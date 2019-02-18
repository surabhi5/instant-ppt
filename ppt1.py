import os
from pptx import Presentation
from bs4 import BeautifulSoup
import requests

finalSummary=[]

def tt(topic):

    response = requests.get("https://en.wikipedia.org/wiki/"+topic)
    soup = BeautifulSoup(response.text,"lxml")
    t = soup.find_all('p') #fetch <p></p> tags 
    
    text = ""
    
    for i in range(0,len(t)):
        #print(t[i].text) #display all text within <p> tags
        text+=t[i].text #add to corpus


    import re
    text=re.sub('\[.*?\]','',text).strip()


##tokenize data into words and sentences
    from nltk.tokenize import word_tokenize, sent_tokenize

##"Today was a great day! How are you? This is a test."
##"Today","was","a",..."!",...
##"Today was a great day! "How are you ?"....

    words=word_tokenize(text)
    
    sentences=sent_tokenize(text)
    from stopwords import stopwords

    _stopwords=word_tokenize(stopwords)

    sentence_scores={}
    word_scores={}







    total_words=len(words) #total no. of words present 


    for word in words:
     if len(word.strip()) > 1 and word.lower not in _stopwords:
            if word.lower() in word_scores.keys():
                word_scores[word.lower()]+=1
            else:
                word_scores[word.lower()]=1 
    for sentence in sentences:
        for word in word_scores.keys():
            if(word in sentence.lower()):
                if sentence in sentence_scores.keys():
                     sentence_scores[sentence]+=word_scores[word]
                else:
                    sentence_scores[sentence]=word_scores[word]
    i=0





#API_KEY="wVjI5flIynUfBCNg0n5HJ9BPgZLaHZM9CW8pAeiyEjI"

#paralleldots.get_api_key()

    #finalSummary=[]
    for sentence in sentence_scores.keys():
        genFact=sentence_scores[sentence]/total_words*3.75
        if(genFact>0.85):
            finalSummary.append(sentence)
            i+=1
           #print(str(i)+">"+sentence+" SCORE : "+str(genFact))
            #print()
            #print()


    
    return(finalSummary)

print("Enter query : ")
query=input() 

current_dir=os.path.dirname(os.path.realpath(__file__))

prs = Presentation(current_dir+"\\abcd.pptx")

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx is here!"


    
result = tt(query)

for sentence_result in result:
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    para=slide.placeholders[1]

    title.text = "explain!"
    para.text=sentence_result

#print(text)

#print(scrape(query))

prs.save('name1.pptx')
