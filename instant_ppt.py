import os
from pptx import Presentation
from bs4 import BeautifulSoup
from pptx.util import Inches
import downloadimage
import requests
from search import scrape
from analyzer import analyze

finalSummary=[]

pos=[] #store positive sentences
neg=[] #store negative sentences
neu=[] #store neutral sentences

def tt(topic):
    url=scrape(topic)
    print("Downloading :" +url[(url.rfind('/')+1):len(url)])
    downloadimage.downloadImage(url[(url.rfind('/')+1):len(url)])
    

    response = requests.get(url)
    soup = BeautifulSoup(response.text,"lxml")
    t = soup.find_all('p') #fetch <p></p> tags 
    
    text = ""
    
    for i in range(0,len(t)):
        #print(t[i].text) #display all text within <p> tags
        text+=t[i].text #add to corpus


    import re
    from summarizer import summarize
    return summarize(text)

    
print("Enter query : ")
query=input() 

result = tt(query)

neu,pos,neg=analyze(result) #analyze sentiment

current_dir=os.path.dirname(os.path.realpath(__file__))

prs = Presentation(current_dir+"\\abcd.pptx")

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = query
subtitle.text = result[0]


img_path = 'image0.jpg'
blank_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(blank_slide_layout)
title=slide.shapes.title
para=slide.placeholders[1]
title.text=" "
para.text=" "


top=Inches(2)
left = Inches(1)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, height=height)


if(len(neu)>0):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "ABOUT"
    subtitle.text = " "


    for s in neu :
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        para=slide.placeholders[1]

        title.text = " "
        para.text=s



img_path = 'image1.jpg'
blank_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(blank_slide_layout)
title=slide.shapes.title
para=slide.placeholders[1]
title.text=" "
para.text=" "


top=Inches(2)
left = Inches(1)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

if(len(pos)>0):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "ADVANTAGES"
    subtitle.text = " "


    for s in pos :
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        para=slide.placeholders[1]

        title.text = " "
        para.text=s




img_path = 'image2.jpg'
blank_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(blank_slide_layout)
title=slide.shapes.title
para=slide.placeholders[1]
title.text=" "
para.text=" "


top=Inches(2)
left = Inches(1)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

if(len(neg)>0):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "DISADVANTAGES"
    subtitle.text = " "

    for s in neg :
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        para=slide.placeholders[1]

        title.text = " "
        para.text=s


prs.save('output.pptx')